#!/usr/local/Caskroom/miniconda/base/bin/python

import os
import pptx
import configparser
import requests
import spacy
import sys

from spacy.lang.en import English

nlp = spacy.load("en_core_web_sm")


def process_file(folder='', api_key='', print_raw=False):
    print('traversing {}'.format(folder))
    if 'clean' in folder:
        return
    for filename in os.listdir(folder):
        file_path = os.path.join(folder, filename)
        if filename.endswith(".pptx"):
            text = get_text_from_ppt(os.path.join(file_path))
            sanitized_string = sanitize_string(text)
            if print_raw is True:
                print(sanitized_string)
            chunks = text_to_chunks(sanitized_string)
            for chunk in chunks:
                ai_response = send_to_open_ai(chunk, api_key)
                print(ai_response)
            print("======== done ===========")
            continue;
        elif os.path.isfile(file_path):
            continue
        else:
            process_file(file_path, api_key)


def send_to_open_ai(text, api_key):
    prompt = 'Give me a synthesize paragraph from the following text: {}'.format(text)
    data = {'prompt': prompt, 'max_tokens': 500}
    headers = {'Authorization': 'Bearer {}'.format(api_key)}
    response = requests.post('https://api.openai.com/v1/engines/text-davinci-003/completions', json=data,
                             headers=headers)
    if (response.status_code == 400):
        return response.text

    raw_response = response.json()
    ai_response = raw_response["choices"][0]["text"]
    return ai_response


def get_text_from_ppt(ppt_path):
    prs = pptx.Presentation(ppt_path)
    print("===== processing file: {} =======".format(ppt_path))
    text = []
    for slide in prs.slides:
        slide_id = prs.slides.index(slide)
        text.append("Title: {} slide: {}".format(ppt_path, slide_id))
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text.append(run.text)
    return '\n'.join(text)


def sanitize_string(s):
    lines = s.split('\n')
    new_lines = []
    for line in lines:
        if len(line.split(' ')) <= 3:
            continue
        if line in new_lines:
            continue
        if line == '':
            continue
        if line.startswith('Title:'):
            continue
        new_lines.append(line)
    return '\n'.join(new_lines)


def text_to_chunks(text):
    chunks = [[]]
    chunk_total_words = 0
    sentences = nlp(text)
    for sentence in sentences.sents:
        chunk_total_words += len(sentence.text.split(" "))
        if chunk_total_words > 1900:
            chunks.append([])
            chunk_total_words = len(sentence.text.split(" "))
        chunks[len(chunks) - 1].append(sentence.text)
    return chunks


def main():
    print_raw = False
    if len(sys.argv) == 2:
        print_raw = True
        print("======== print raw input parameter is set.. printing in console ppt contents")
    config = configparser.ConfigParser()
    config.read(os.path.join(os.path.dirname(__file__), 'properties.ini'))
    directory_ = config['DEFAULT']['Local_Directory']
    print("reading: {}".format(directory_))
    process_file(directory_, config['DEFAULT']['API_KEY'], print_raw)


if __name__ == "__main__":
    main()
